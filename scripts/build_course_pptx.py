#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成 MusicEditing 硬核课 PPT（多期）。

用法（仓库根）:
  python scripts/build_course_pptx.py

输出: docs/course/pptx/*.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "course" / "pptx"

BG = RGBColor(0x12, 0x14, 0x18)
FG = RGBColor(0xE8, 0xED, 0xF5)
ACCENT = RGBColor(0xE8, 0xA4, 0x5C)
MUTED = RGBColor(0x8B, 0x95, 0xA8)


def _set_run(run, text: str, *, size: int = 20, bold: bool = False, color=FG) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"


def _fill_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    tbox = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.8))
    tr = tbox.text_frame.paragraphs[0].add_run()
    _set_run(tr, title, size=34, bold=True, color=ACCENT)
    sbox = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(12), Inches(1.4))
    sbox.text_frame.word_wrap = True
    sr = sbox.text_frame.paragraphs[0].add_run()
    _set_run(sr, subtitle, size=18, color=MUTED)
    bbox = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(12), Inches(0.4))
    br = bbox.text_frame.paragraphs[0].add_run()
    _set_run(br, "MusicEditing · 硬核实战课 · docs/course", size=12, color=MUTED)


def add_section(prs: Presentation, title: str, bullets: list[str], *, footer: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    head = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(0.8))
    _set_run(head.text_frame.paragraphs[0].add_run(), title, size=28, bold=True, color=ACCENT)

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        text = line.strip()
        if not text.startswith("•"):
            text = "• " + text
        _set_run(p.add_run(), text, size=18, color=FG)

    if footer:
        ft = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12), Inches(0.35))
        _set_run(ft.text_frame.paragraphs[0].add_run(), footer, size=11, color=MUTED)


def build_deck(code: str, title: str, subtitle: str, sections: list[tuple[str, list[str]]]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs, f"{code}  {title}", subtitle)
    for i, (sec_title, bullets) in enumerate(sections, 1):
        add_section(prs, sec_title, bullets, footer=f"{code} · {i}/{len(sections)} · MusicEditing")
    OUT.mkdir(parents=True, exist_ok=True)
    # 文件名用 ASCII，避免部分环境下中文路径乱码
    path = OUT / f"{code}.pptx"
    prs.save(path)
    return path


DECKS: list[tuple[str, str, str, list[tuple[str, list[str]]]]] = [
    (
        "00",
        "课程总览",
        "12 期 × 90–120 分钟 · 可跑仓库 · 可改代码 · 可打包上线",
        [
            ("这门课卖的是什么", [
                "不是剪映拖拽入门，不是空洞架构动画",
                "交付：本地 AI 音视频工具从引擎到发版的完整工程能力",
                "技术栈：C++ FFmpeg + PySide6 MVVM + ONNX + 便携分发 + 卡密商业化",
                "期末标准：能讲清调用链、试用门禁、更新/签名/收银台边界",
            ]),
            ("四模块路线", [
                "A 地基 S01–S02：产品全景 + x64 硬核跑通",
                "B 架构 S03–S05：MVVM/UI、media_engine/CLI、播放器 IPC/SHM",
                "C 业务 S06–S10：切片、超分、去水印/队列、下载热评、趣味扩展",
                "D 商业 S11–S12：试用卡密、打包 Inno、更新通道、验收毕业",
            ]),
            ("学习资产", [
                "PPT：docs/course/pptx/",
                "讲师讲稿：docs/course/sessions/",
                "自学路径：docs/LEARNING.md",
                "实现真源：docs/design/（上课引用，不替代课件）",
            ]),
            ("课堂纪律", [
                "archive/ 与本机 cookies/log 不是教材",
                "改功能必须同步 design 专文（见 Cursor skill）",
                "每期有 Lab；S12 毕业作业必须 accept PASS",
            ]),
        ],
    ),
    (
        "S01",
        "开班与产品全景",
        "为什么做成「能卖的本地工具」而不是 Demo",
        [
            ("开场：真实痛点", [
                "隐私与素材不出网：本地处理是卖点不是情怀",
                "平台角标/水印、批量无人值守、高光切片：刚需场景",
                "云剪辑贵、限速、条款风险：对照本地交付边界",
            ]),
            ("产品边界（讲课必说）", [
                "纯本地离线优先；联网仅更新/激活/下载可选",
                "不是专业 NLE 多轨时间线（明确不做的大坑）",
                "功能以菜单为准：核心 / 工作流 / 趣味 / 帮助",
            ]),
            ("功能全景速览", [
                "核心：首页预览、智能切片、画质增强、去水印",
                "工作流：全流程队列、素材库、下载热评、BGM",
                "趣味：热评弹幕、封面工厂、音频梗音、溯源水印",
                "帮助：个人中心（授权/GPU/诊断/更新）",
            ]),
            ("架构一句话", [
                "View(PySide6) ↔ MainViewModel ↔ MediaBridge/PlayerBackend",
                "→ media_cli / media_player → media_engine.dll（FFmpeg/OpenCV/ORT）",
                "短调用：ctypes 直连 DLL，失败回退 CLI；播放走子进程",
            ]),
            ("文档真源纪律", [
                "实现真源：docs/design/",
                "自学路径：docs/LEARNING.md",
                "产品交互稿：只读对照；冲突以 design 为准",
                "docs/archive/：历史科普，禁止当教材",
            ]),
            ("12 期地图与期末标准", [
                "S02 跑通 → S03–S05 架构 → S06–S10 业务 → S11–S12 商业发版",
                "期末：口述调用链 + 试用/正式差异 + 四件自备项",
                "Lab：画「菜单 → 用户价值」图（下期带来）",
            ]),
        ],
    ),
    (
        "S02",
        "仓库与硬核跑通",
        "clone → setup_ffmpeg → build_x64 → run_ui_x64",
        [
            ("环境清单", [
                "Windows 10/11 x64 + VS 2022/2026（C++ 桌面）",
                "CMake 3.20+ · Python 3.10+ · PySide6",
                "日常只推 x64；Win32 并存但扩展以 x64 为主",
            ]),
            ("命令链（当场演示）", [
                "scripts\\setup_ffmpeg_x64.bat",
                "build_x64.bat",
                "pip install -r client\\scripts\\requirements.txt",
                "run_ui_x64.bat",
            ]),
            ("产物必须认识", [
                "build_x64\\bin\\Release\\media_cli.exe — 批处理入口",
                "media_player.exe — 首页播放子进程",
                "media_engine.dll — C API（probe/thumbnail/超分…）",
                "配置：client/resources/config/app.conf",
            ]),
            ("可选依赖怎么讲", [
                "OpenCV / ORT / 超分·LaMa·Vosk 模型：缺了 UI 可开，功能提示",
                "models/README.md + scripts\\download_*.bat",
                "不要把「没下模型」说成「项目坏了」",
            ]),
            ("常见翻车", [
                "没跑 setup_ffmpeg_x64 / 架构混用 Win32 产物",
                "VS 未装 C++ 工作负载导致引擎链失败",
                "pip 装到错误解释器；应用 64-bit Python",
            ]),
            ("Lab 四验收（截图交作业）", [
                "UI 窗口已打开",
                "首页能播 tests\\test_video.mp4",
                "media_cli.exe 存在",
                "media_player.exe 存在",
            ]),
        ],
    ),
    (
        "S03",
        "MVVM与Studio_UI",
        "View / ViewModel / Model · 懒加载 · 菜单索引",
        [
            ("MVVM 在本仓库的落地", [
                "View：client/scripts/ui/*.py",
                "ViewModel：viewmodels/main_vm.py（Signal/Slot 跨线程）",
                "Model：dataclass；重活在后台线程，UI 只收 Signal",
            ]),
            ("导航硬约束", [
                "workflow_link.TAB_* 索引稳定，勿随意改号",
                "MENU_GROUPS：核心/工作流/趣味/帮助",
                "QStackedWidget 懒创建页 + 空闲预热减首点卡顿",
            ]),
            ("Studio 体验细节（值钱点）", [
                "studio_kit：Hero/Card/wrap_studio_scroll/wrap_tab_scroll",
                "ElidedPathLabel：长路径不撑爆布局",
                "GroupBox 标题留白，避免字压控件",
                "封面工厂等长页：滚动 + 底栏按钮",
            ]),
            ("接力与完成弹窗", [
                "open_with_video / ask_video_handoff",
                "导出完成可「打开文件夹」选中成片",
                "三大功能串联 ≠ 队列无人值守（下几期展开）",
            ]),
            ("Lab", [
                "改一处 UI 文案或 app.conf 注释",
                "对照 music-editing-feature-docs skill：要不要改 feature_flows/状态表",
                "说明：纯文案通常不动状态表；行为变更必须同步文档",
            ]),
        ],
    ),
    (
        "S04",
        "引擎与CLI协议",
        "media_engine.dll · media_cli · ctypes 优先",
        [
            ("三件套分工", [
                "DLL：探测/遍历/缩略图/超分·去水印 C API",
                "CLI：Python 批处理入口；stdout=协议 stderr=日志",
                "Player：仅首页预览子进程，不与 CLI 混用",
            ]),
            ("MediaBridge 策略", [
                "probe/thumbnail：优先 ctypes 直连 DLL",
                "失败回退 media_cli.exe",
                "mtime 缓存减少重复探测",
            ]),
            ("导出与 FFmpeg", [
                "高光分段/拼接优先 -c copy remux",
                "需要时整段重编码；AAC + faststart",
                "捆绑 ffmpeg 保证别人电脑可复现",
            ]),
            ("协议课怎么上硬核", [
                "打开 media_engine.md 的 CLI 示例",
                "强调：UI 不解析「人类日志」，只认 stdout 协议行",
                "改 CLI 输出格式 = 破坏兼容，必须改文档与回归",
            ]),
            ("Lab", [
                "按文档跑一次 probe 或 thumbnail",
                "把关键 stdout 贴进作业（可打码路径）",
            ]),
        ],
    ),
    (
        "S05",
        "播放器IPC与SHM",
        "子进程 · 双缓冲 · 异步 Seek",
        [
            ("一帧上屏路径", [
                "VideoPlayerWidget → PlayerBackend",
                "stdin/stdout 控制 media_player.exe",
                "解码 RGB →（SHM 双缓冲）→ OpenGL/显示控件",
                "音频走 Qt 侧，音画软校正",
            ]),
            ("为什么要 SHM / 预取", [
                "避免大帧走管道拷贝",
                "双缓冲减少撕裂与等待",
                "lookahead 预取降低 Seek 后空窗",
            ]),
            ("Seek 硬核细节", [
                "异步首帧：UI 不卡死在 seek_and_frame",
                "Seek 期间抑制软校正抢状态",
                "松手后预热预取再恢复播放",
            ]),
            ("硬解与回退", [
                "可选 D3D11VA；失败回退 CPU",
                "滤镜 OpenCL UMat 可关；稳定性优先",
            ]),
            ("Lab", [
                "对照流程图 README + player_decode_flow",
                "手绘/白板：从点击播放到像素出现",
            ]),
        ],
    ),
    (
        "S06",
        "智能切片与高光",
        "演讲 · 游戏 · 响度 · 竖屏成片",
        [
            ("场景分流", [
                "演讲金句：Vosk ASR + LLM/规则；无人声模型则人声段兜底",
                "游戏高光：切点 + 运动/闪光；game_event.onnx 现为 stub",
                "响度高潮：ebur128；手动切片：不依赖模型",
            ]),
            ("时间轴体验", [
                "缩略图胶片条 + 列表图标 §5.1.1",
                "后台抽中点缩略图，不堵 UI",
            ]),
            ("成片能力", [
                "一键高光成片 · 静音剪掉",
                "竖屏 9:16：锚点 / 智能跟脸",
                "发布预设与规范命名（抖音/B站/快手）",
            ]),
            ("试用门禁埋伏笔", [
                "高光导出次数 · 竖屏次数 · 最长边≤720p",
                "正式版解锁在 S11 展开",
            ]),
            ("Lab", [
                "手动两段高光 → 竖屏导出",
                "记录是否触发试用限制",
            ]),
        ],
    ),
    (
        "S07",
        "画质增强硬核",
        "OpenCV vs Real-ESRGAN · tile · 补帧 · LUT",
        [
            ("两条超分路径", [
                "快速：OpenCV（JPEG 中间帧 + 多线程）",
                "AI：Real-ESRGAN ONNX（Session 串行）",
                "对比预览：左右、滚轮缩放、拖拽平移",
            ]),
            ("性能旋钮", [
                "CUDA EP 探测与缓存；无 EP 明示走 CPU",
                "自动 tile≈640（CUDA）/ 更小省显存",
                "视频默认试跑秒数：批量可感、可演示",
            ]),
            ("补帧与调色", [
                "minterpolate 快速/精细；可选 RIFE",
                "LUT / 一键调色与 FrameProcessor 预设对齐",
            ]),
            ("授权差异", [
                "试用：OpenCV 2× 可用",
                "正式：AI 4× / 相关能力解锁",
            ]),
            ("Lab", [
                "短视频 OpenCV 2× 试跑，记录耗时与输出",
            ]),
        ],
    ),
    (
        "S08",
        "去水印与全流程队列",
        "角标 · 批量重试 · 有限并行",
        [
            ("去水印路径", [
                "视频默认快速 OpenCV；图片/精修 LaMa",
                "框选多区域；平台角标预设（抖/快）",
                "批量：失败重试、结果列表",
            ]),
            ("队列编排（高含金量）", [
                "步骤可勾选：切片 → 超分 → 去水印",
                "max_parallel：切片/导出可重叠",
                "超分+去水印信号量串行：防 GPU/磁盘互抢",
                "成片模板 · 产物 GB 上限 · 分阶段 ETA",
            ]),
            ("与「完成弹窗接力」的区别", [
                "单页完成：弹窗送去下一站",
                "队列：无人值守，不走各页弹窗",
            ]),
            ("Lab", [
                "队列只开切片 + OpenCV 超分试跑 8s",
                "跑通 tests 视频并保留输出目录截图",
            ]),
        ],
    ),
    (
        "S09",
        "下载热评与成片模板",
        "yt-dlp Cookie · ASS · 发布预设",
        [
            ("下载硬核现实", [
                "链接格式通常没问题；站点要 Cookie",
                "推荐 Netscape cookies.txt 绝对路径",
                "浏览器 DPAPI/锁库常失败：讲清原因",
            ]),
            ("失败体验也是产品", [
                "限流 / 无音轨 / Cookie 失效：白话提示",
                "换 Cookie / 重试路径写进 UI",
            ]),
            ("热评三合一", [
                "拉取 · 滚动 · 成片（danmaku/cards/ass）",
                "送首页播放叠弹幕：速度/密度/区域",
            ]),
            ("发布向", [
                "竖屏模板 + 封面话题草稿",
                "规范命名与导出参数面板",
            ]),
            ("Lab", [
                "配置 yt_dlp_cookies_file（勿提交 git）",
                "书面解释：为何 cookies 不能进仓库",
            ]),
        ],
    ),
    (
        "S10",
        "趣味扩展包",
        "封面 · 音频梗音 · BGM · 溯源水印",
        [
            ("封面工厂", [
                "均匀抽样 + Laplacian 选清晰帧",
                "大标题 PNG；可选 EXIF / 频域溯源",
            ]),
            ("音频趣味 / 梗音", [
                "整轨：变调变速倒放伪8D混响",
                "梗音叠加：用户自备热梗，注意版权声明",
            ]),
            ("BGM 与 Demucs", [
                "基础混音仅 FFmpeg，可打包分发",
                "Demucs 可选：PyTorch 体积大，setup 脚本单独装",
            ]),
            ("溯源水印边界", [
                "频域 / 回声 / LSB / EXIF",
                "不宣称对抗平台重编码——诚信卖点",
            ]),
            ("Lab", [
                "生成封面 PNG；可选打开溯源勾选验证",
            ]),
        ],
    ),
    (
        "S11",
        "商业闭环",
        "trial_policy · 卡密 · 激活服 · 购买页",
        [
            ("试用策略", [
                "次数：高光 / 竖屏配额",
                "质量：最长边 ≤720p",
                "能力：AI4× / 队列 / LaMa 等门禁",
            ]),
            ("卡密与激活", [
                "本地格式校验 + license_fp",
                "可选 POST /v1/activate 联网激活",
                "scripts/license_server：签发 + 演示购买页",
            ]),
            ("钱在店外", [
                "微信/支付宝/Lemon 等收银台不进本仓库",
                "支付成功 → 发卡 → 用户在个人中心兑换",
                "app.conf：license_purchase_url / license_server_url",
            ]),
            ("Lab", [
                "启动 license_server，走一遍演示激活",
                "截图：试用 → 正式 →（可选）恢复试用",
            ]),
        ],
    ),
    (
        "S12",
        "发版毕业课",
        "pack · accept · Inno · 更新 · 签名自备项",
        [
            ("便携包档位", [
                "slim / standard / full",
                "默认 embed Python + .pyc；可选 --sign",
                "accept_portable：机器预检；干净机仍要人手",
            ]),
            ("安装与一键发版", [
                "Inno：build_installer.bat",
                "release_oneclick：回归→pack→accept→Inno→清单",
            ]),
            ("自动更新通道", [
                "publish_update_manifest → dist/update/",
                "serve_update_channel 仅本地联调",
                "发版包禁止残留 127.0.0.1 update URL",
                "启动静默检查 + update_last_notified 防刷",
            ]),
            ("四件必须自备（背诵）", [
                "干净机手测（SmartScreen / VC++）",
                "CDN 静态托管更新目录",
                "代码签名证书 thumbprint",
                "外部收银台与发卡流程",
            ]),
            ("毕业 Lab", [
                "pack slim --zip",
                "accept_portable → PASS",
                "publish_update_manifest 生成 json",
                "提交一页上线检查表（含自备项）",
            ]),
            ("结营", [
                "回看调用链：View→VM→Bridge→CLI/Player→DLL",
                "你卖的是可维护的本地工具工程能力",
                "持续以 docs/design 为真源迭代",
            ]),
        ],
    ),
]


def main() -> int:
    print("生成课件…")
    for code, title, sub, sections in DECKS:
        path = build_deck(code, title, sub, sections)
        print(f"  OK {path.name}")
    print(f"输出目录: {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
